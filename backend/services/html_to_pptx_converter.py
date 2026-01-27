"""
HTML 到 PPTX 转换器 - 修复重复文字问题

核心策略：
1. 先提取文本元素的位置和样式信息
2. 隐藏所有文本元素
3. 截图背景（不含文字）
4. 在 PPTX 中添加可编辑的文本框

这样就避免了文字重复显示的问题！

作者：PPTAgent Team
日期：2026-01-21
"""

import logging
import tempfile
import shutil
from typing import List, Dict, Any, Optional
from pathlib import Path

from playwright.async_api import async_playwright, Browser, Page
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)


class HTMLToPPTXConverter:
    """HTML 到 PPTX 转换器 - 避免文字重复"""
    
    SLIDE_WIDTH_PX = 1920
    SLIDE_HEIGHT_PX = 1080
    SLIDE_WIDTH_INCH = 10
    SLIDE_HEIGHT_INCH = 5.625
    DPI = 96
    
    def __init__(self):
        self.temp_dir: Optional[Path] = None
        self.slides_info: List[Dict] = []
        
    async def convert(
        self,
        slides_html: List[str],
        output_path: str,
        title: str = "Presentation"
    ) -> str:
        """转换 HTML 幻灯片为 PPTX 文件"""
        self.temp_dir = Path(tempfile.mkdtemp(prefix="html2pptx_"))
        self.slides_info = []
        
        try:
            logger.info(f"开始转换 {len(slides_html)} 张幻灯片...")
            
            # 保存 HTML 文件
            for i, html in enumerate(slides_html):
                html_file = self.temp_dir / f"slide_{i+1:02d}.html"
                
                if not html.strip().startswith('<!DOCTYPE') and not html.strip().startswith('<html'):
                    html = self._wrap_html(html)
                
                html_file.write_text(html, encoding='utf-8')
                self.slides_info.append({
                    'slide_number': i + 1,
                    'html_path': str(html_file)
                })
            
            # 使用 Playwright 分析
            async with async_playwright() as p:
                browser = await p.chromium.launch(
                    headless=True,
                    args=['--no-sandbox', '--disable-setuid-sandbox']
                )
                
                try:
                    all_analyses = await self._analyze_all_slides(browser)
                finally:
                    await browser.close()
            
            # 构建 PPTX
            presentation = self._create_presentation()
            
            for i, analysis in enumerate(all_analyses, 1):
                try:
                    if 'error' in analysis:
                        self._create_error_slide(presentation, analysis['error'], i)
                    else:
                        self._build_slide_from_analysis(presentation, analysis)
                        logger.info(f"幻灯片 {i} 构建完成")
                except Exception as e:
                    logger.error(f"构建幻灯片 {i} 失败: {e}")
                    self._create_error_slide(presentation, str(e), i)
            
            presentation.core_properties.title = title
            presentation.core_properties.author = "PPTAgent"
            
            presentation.save(output_path)
            logger.info(f"PPTX 文件已保存: {output_path}")
            
            return output_path
            
        finally:
            if self.temp_dir and self.temp_dir.exists():
                shutil.rmtree(self.temp_dir)
    
    def _wrap_html(self, html: str) -> str:
        """包装 HTML 为完整文档"""
        return f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <style>
        * {{ box-sizing: border-box; }}
        body {{ 
            margin: 0; 
            padding: 0; 
            width: {self.SLIDE_WIDTH_PX}px; 
            height: {self.SLIDE_HEIGHT_PX}px;
            overflow: hidden;
        }}
    </style>
</head>
<body>
{html}
</body>
</html>'''
    
    def _create_presentation(self) -> Presentation:
        """创建演示文稿"""
        prs = Presentation()
        prs.slide_width = Inches(self.SLIDE_WIDTH_INCH)
        prs.slide_height = Inches(self.SLIDE_HEIGHT_INCH)
        
        if len(prs.slides) > 0:
            xml_slides = prs.slides._sldIdLst
            for i in range(len(xml_slides) - 1, -1, -1):
                xml_slides.remove(xml_slides[i])
        
        return prs
    
    def _create_error_slide(self, presentation: Presentation, error: str, slide_num: int):
        """创建错误占位幻灯片"""
        blank_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(blank_layout)
        
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(2)
        )
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = f"错误：幻灯片 {slide_num} 处理失败\n{error}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 0, 0)
    
    async def _analyze_all_slides(self, browser: Browser) -> List[Dict]:
        """分析所有幻灯片"""
        all_analyses = []
        
        for slide_info in self.slides_info:
            try:
                logger.info(f"分析幻灯片 {slide_info['slide_number']}...")
                
                page = await browser.new_page(
                    viewport={
                        'width': self.SLIDE_WIDTH_PX, 
                        'height': self.SLIDE_HEIGHT_PX
                    }
                )
                
                try:
                    # 加载页面
                    await page.goto(
                        f"file://{slide_info['html_path']}", 
                        wait_until='networkidle',
                        timeout=60000
                    )
                    await page.wait_for_timeout(1000)
                    
                    # ⭐ 关键步骤1：先提取文本信息（此时文本还可见）
                    text_elements = await self._extract_text_elements(page)
                    
                    # ⭐ 关键步骤2：隐藏所有文本元素
                    await self._hide_text_elements(page)
                    
                    # ⭐ 关键步骤3：截图背景（不含文字）
                    bg_path = self.temp_dir / f"slide_{slide_info['slide_number']}_bg.png"
                    await page.screenshot(path=str(bg_path), full_page=False)
                    
                    # 步骤4：提取视觉元素（图片、Canvas）
                    visual_elements = await self._extract_visual_elements(
                        page, slide_info['slide_number']
                    )
                    
                    analysis = {
                        'slide_info': slide_info,
                        'background_path': bg_path,
                        'text_elements': text_elements,
                        'visual_elements': visual_elements
                    }
                    
                    all_analyses.append(analysis)
                    logger.info(
                        f"幻灯片 {slide_info['slide_number']} 分析完成: "
                        f"{len(text_elements)} 文本, {len(visual_elements)} 视觉元素"
                    )
                    
                finally:
                    await page.close()
                    
            except Exception as e:
                logger.error(f"分析幻灯片 {slide_info['slide_number']} 失败: {e}")
                all_analyses.append({
                    'slide_info': slide_info,
                    'error': str(e)
                })
        
        return all_analyses
    
    async def _hide_text_elements(self, page: Page):
        """隐藏所有文本元素 - 关键函数！"""
        await page.evaluate('''() => {
            // 找到所有包含文本的元素
            function hasDirectTextContent(el) {
                for (const node of el.childNodes) {
                    if (node.nodeType === Node.TEXT_NODE) {
                        const text = node.textContent.trim();
                        if (text.length > 0) return true;
                    }
                }
                return false;
            }
            
            const allElements = document.body.querySelectorAll('*');
            for (const el of allElements) {
                // 跳过这些元素
                const skipTags = ['SCRIPT', 'STYLE', 'IMG', 'SVG', 'CANVAS', 'VIDEO', 'IFRAME'];
                if (skipTags.includes(el.tagName)) continue;
                
                if (hasDirectTextContent(el)) {
                    // 隐藏文本但保留布局
                    el.style.color = 'transparent';
                    // 或者使用：el.style.visibility = 'hidden';
                    // 或者使用：el.style.opacity = '0';
                }
            }
        }''')
    
    async def _extract_text_elements(self, page: Page) -> List[Dict]:
        """提取文本元素信息"""
        text_elements = await page.evaluate('''() => {
            const elements = [];
            const DPI = 96;
            
            function pxToInch(px) {
                return px / DPI;
            }
            
            function rgbToHex(rgb) {
                if (!rgb || rgb === 'transparent' || rgb === 'rgba(0, 0, 0, 0)') {
                    return '#000000';
                }
                const match = rgb.match(/rgba?\\((\\d+),\\s*(\\d+),\\s*(\\d+)/);
                if (!match) return '#000000';
                return '#' + match.slice(1, 4)
                    .map(n => parseInt(n).toString(16).padStart(2, '0'))
                    .join('')
                    .toUpperCase();
            }
            
            function isVisible(el) {
                const rect = el.getBoundingClientRect();
                if (rect.width <= 1 || rect.height <= 1) return false;
                
                const style = window.getComputedStyle(el);
                if (style.display === 'none') return false;
                if (style.visibility === 'hidden') return false;
                if (parseFloat(style.opacity) < 0.1) return false;
                
                return true;
            }
            
            function hasDirectTextContent(el) {
                for (const node of el.childNodes) {
                    if (node.nodeType === Node.TEXT_NODE) {
                        const text = node.textContent.trim();
                        if (text.length > 0) return true;
                    }
                }
                return false;
            }
            
            function shouldExtractText(el) {
                const skipTags = ['SCRIPT', 'STYLE', 'IMG', 'SVG', 'CANVAS', 'VIDEO', 'IFRAME'];
                if (skipTags.includes(el.tagName)) return false;
                
                if (!isVisible(el)) return false;
                if (!hasDirectTextContent(el)) return false;
                
                // 检查是否有块级子元素
                const children = Array.from(el.children);
                const hasBlockChildren = children.some(child => {
                    const style = window.getComputedStyle(child);
                    return style.display === 'block' || 
                           style.display === 'flex' ||
                           style.display === 'grid';
                });
                
                return !hasBlockChildren;
            }
            
            const allElements = document.body.querySelectorAll('*');
            const processedRects = new Set();
            
            for (const el of allElements) {
                if (!shouldExtractText(el)) continue;
                
                const rect = el.getBoundingClientRect();
                const rectKey = `${rect.left.toFixed(1)},${rect.top.toFixed(1)}`;
                
                if (processedRects.has(rectKey)) continue;
                processedRects.add(rectKey);
                
                const style = window.getComputedStyle(el);
                const text = el.innerText?.trim() || '';
                
                if (text.length > 0) {
                    const fontSize = parseFloat(style.fontSize) * 0.75;
                    
                    elements.push({
                        type: 'text',
                        content: text,
                        x: pxToInch(rect.left),
                        y: pxToInch(rect.top),
                        width: pxToInch(rect.width),
                        height: pxToInch(rect.height),
                        font_size: Math.max(8, fontSize),
                        font_family: style.fontFamily.split(',')[0]
                            .replace(/['"]/g, '').trim() || 'Arial',
                        color: rgbToHex(style.color),
                        bold: parseInt(style.fontWeight) >= 600,
                        italic: style.fontStyle === 'italic',
                        underline: style.textDecoration.includes('underline'),
                        align: style.textAlign === 'start' ? 'left' : 
                               (style.textAlign || 'left'),
                        z_index: parseInt(style.zIndex) || 0
                    });
                }
            }
            
            elements.sort((a, b) => {
                if (a.z_index !== b.z_index) return a.z_index - b.z_index;
                if (Math.abs(a.y - b.y) > 0.1) return a.y - b.y;
                return a.x - b.x;
            });
            
            return elements;
        }''')
        
        return text_elements
    
    async def _extract_visual_elements(self, page: Page, slide_number: int) -> List[Dict]:
        """提取视觉元素（图片、Canvas、SVG）"""
        visual_elements = []
        
        # 提取图片
        images = await page.query_selector_all('img[src]')
        for i, img in enumerate(images):
            try:
                is_visible = await img.is_visible()
                if not is_visible:
                    continue
                
                rect = await img.bounding_box()
                if rect and rect['width'] > 10 and rect['height'] > 10:
                    img_path = self.temp_dir / f"slide_{slide_number}_img_{i+1}.png"
                    await img.screenshot(path=str(img_path))
                    
                    visual_elements.append({
                        'type': 'image',
                        'x': rect['x'] / self.DPI,
                        'y': rect['y'] / self.DPI,
                        'width': rect['width'] / self.DPI,
                        'height': rect['height'] / self.DPI,
                        'image_path': img_path,
                        'z_index': 1
                    })
            except Exception as e:
                logger.warning(f"提取图片 {i+1} 失败: {e}")
        
        # 提取 Canvas（图表等）
        canvases = await page.query_selector_all('canvas')
        for i, canvas in enumerate(canvases):
            try:
                is_visible = await canvas.is_visible()
                if not is_visible:
                    continue
                
                rect = await canvas.bounding_box()
                if rect and rect['width'] > 10 and rect['height'] > 10:
                    canvas_path = self.temp_dir / f"slide_{slide_number}_canvas_{i+1}.png"
                    await canvas.screenshot(path=str(canvas_path))
                    
                    visual_elements.append({
                        'type': 'canvas',
                        'x': rect['x'] / self.DPI,
                        'y': rect['y'] / self.DPI,
                        'width': rect['width'] / self.DPI,
                        'height': rect['height'] / self.DPI,
                        'image_path': canvas_path,
                        'z_index': 1
                    })
            except Exception as e:
                logger.warning(f"提取 Canvas {i+1} 失败: {e}")
        
        # 提取 SVG
        svgs = await page.query_selector_all('svg')
        for i, svg in enumerate(svgs):
            try:
                is_visible = await svg.is_visible()
                if not is_visible:
                    continue
                
                rect = await svg.bounding_box()
                if rect and rect['width'] > 10 and rect['height'] > 10:
                    svg_path = self.temp_dir / f"slide_{slide_number}_svg_{i+1}.png"
                    await svg.screenshot(path=str(svg_path))
                    
                    visual_elements.append({
                        'type': 'svg',
                        'x': rect['x'] / self.DPI,
                        'y': rect['y'] / self.DPI,
                        'width': rect['width'] / self.DPI,
                        'height': rect['height'] / self.DPI,
                        'image_path': svg_path,
                        'z_index': 1
                    })
            except Exception as e:
                logger.warning(f"提取 SVG {i+1} 失败: {e}")
        
        return visual_elements
    
    def _build_slide_from_analysis(self, presentation: Presentation, analysis: Dict):
        """从分析结果构建幻灯片"""
        background_path = analysis['background_path']
        text_elements = analysis['text_elements']
        visual_elements = analysis['visual_elements']
        
        blank_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(blank_layout)
        
        # 1. 添加背景（已经隐藏了文字）
        if background_path and background_path.exists():
            try:
                slide.shapes.add_picture(
                    str(background_path),
                    Inches(0),
                    Inches(0),
                    width=Inches(self.SLIDE_WIDTH_INCH),
                    height=Inches(self.SLIDE_HEIGHT_INCH)
                )
            except Exception as e:
                logger.warning(f"添加背景失败: {e}")
        
        # 2. 添加视觉元素（其实背景已经包含了，这里可选）
        # 如果视觉元素已经在背景截图里了，这部分可以注释掉
        for visual_elem in visual_elements:
            try:
                if visual_elem['image_path'].exists():
                    slide.shapes.add_picture(
                        str(visual_elem['image_path']),
                        Inches(visual_elem['x']),
                        Inches(visual_elem['y']),
                        width=Inches(visual_elem['width']),
                        height=Inches(visual_elem['height'])
                    )
            except Exception as e:
                logger.warning(f"添加视觉元素失败: {e}")
        
        # 3. 添加可编辑的文本框
        for text_elem in text_elements:
            try:
                self._create_transparent_textbox(slide, text_elem)
            except Exception as e:
                logger.warning(f"添加文本框失败: {e}")
    
    def _create_transparent_textbox(self, slide, text_elem: Dict):
        """创建透明背景的文本框"""
        left = max(0, Inches(text_elem['x']))
        top = max(0, Inches(text_elem['y']))
        width = min(Inches(text_elem['width']), Inches(self.SLIDE_WIDTH_INCH) - left)
        height = min(Inches(text_elem['height']), Inches(self.SLIDE_HEIGHT_INCH) - top)
        
        if width <= 0 or height <= 0:
            return
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        
        # ⭐ 关键：设置完全透明
        textbox.fill.background()
        textbox.line.fill.background()
        
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.01)
        tf.margin_bottom = Inches(0.01)
        
        tf.clear()
        p = tf.paragraphs[0]
        
        # 对齐
        align_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        p.alignment = align_map.get(text_elem.get('align', 'left'), PP_ALIGN.LEFT)
        
        # 添加文本
        run = p.add_run()
        run.text = text_elem['content']
        
        # 字体样式
        font = run.font
        font.size = Pt(text_elem.get('font_size', 12))
        
        font_family = text_elem.get('font_family', '').strip()
        if font_family and font_family.lower() not in ['sans-serif', 'serif', 'monospace']:
            font.name = font_family
        else:
            font.name = 'Microsoft YaHei'
        
        # 颜色
        if text_elem.get('color'):
            try:
                color_hex = text_elem['color'].lstrip('#')
                if len(color_hex) == 6:
                    r = int(color_hex[0:2], 16)
                    g = int(color_hex[2:4], 16)
                    b = int(color_hex[4:6], 16)
                    font.color.rgb = RGBColor(r, g, b)
            except Exception:
                pass
        
        if text_elem.get('bold'):
            font.bold = True
        if text_elem.get('italic'):
            font.italic = True
        if text_elem.get('underline'):
            font.underline = True


# 便捷函数
async def convert_html_to_pptx(
    slides_html: List[str],
    output_path: str,
    title: str = "Presentation"
) -> str:
    """
    转换 HTML 幻灯片为 PPTX 文件
    
    工作流程：
    1. 提取文本信息（位置、样式、内容）
    2. 隐藏所有文本元素
    3. 截图背景（不含文字）
    4. 在 PPTX 中重新添加可编辑的文本框
    
    Args:
        slides_html: HTML 幻灯片列表
        output_path: 输出文件路径
        title: 演示文稿标题
        
    Returns:
        生成的 PPTX 文件路径
    """
    converter = HTMLToPPTXConverter()
    return await converter.convert(slides_html, output_path, title)