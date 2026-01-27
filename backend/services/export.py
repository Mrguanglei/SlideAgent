"""
PPTAgent 导出服务模块

支持将 PPT 导出为：
- PDF 文件
- PNG 图片（打包为 ZIP）
- PPTX 文件（使用 python-pptx 生成可编辑的 PPTX）

更新日志：
- 2026-01-17: 重构 PPTX 导出
  - 使用 Python + python-pptx 替代 Node.js html2pptx
  - 通过 HTML 解析生成可编辑的 PPTX（文本、形状、图片）
  - 解决了 PPTX 文件损坏问题
  - 简化了依赖（不再需要 Node.js 和 pptxgenjs）
- 2026-01-17: 添加导出优化功能
  - 添加 HTML 验证和导出前预检
  - 添加导出质量报告
  - 优化 Playwright 等待时间
  - 优化字体处理
"""

import os
import io
import uuid
import zipfile
import logging
import asyncio
from typing import List, Optional, Tuple, Dict, Any
from pathlib import Path
from datetime import datetime



# 图片生成
from playwright.async_api import async_playwright

# PPTX 生成（新方案：HTML 渲染 + 元素截图 + python-pptx）
from services.html_to_pptx_converter import convert_html_to_pptx

# HTML 验证
from services.html_validator import validate_html

# HTML 静态化
from services.html_staticizer import batch_staticize_html

# HTML 解析
from bs4 import BeautifulSoup
import re

logger = logging.getLogger(__name__)

# 导出文件存储目录
EXPORT_DIR = Path("/tmp/ppt_exports")
EXPORT_DIR.mkdir(parents=True, exist_ok=True)


class PPTExporter:
    """PPT 导出器"""
    
    FONT_REPLACEMENTS = {
        "MiSans": "Microsoft YaHei",
        "Noto Sans SC": "Microsoft YaHei",
        "Source Han Serif SC": "SimSun",
        "Roboto Flex": "Arial",
        "Source Code Pro": "Courier New",
        "抖音黑体": "Microsoft YaHei",
    }
    
    def __init__(self):
        self.slide_width = 1280
        self.slide_height = 720
    
    def _parse_slide_html(self, html: str) -> Tuple[List[str], str]:
        """
        解析幻灯片 HTML，提取样式和 body 内容
        
        Returns:
            (styles, body_content)
        """
        soup = BeautifulSoup(html, 'html.parser')
        
        # 提取所有 <style> 标签
        styles = []
        for style_tag in soup.find_all('style'):
            styles.append(style_tag.string or "")
        
        # 提取 <body> 内容
        body = soup.find('body')
        if body:
            # 移除 <body> 标签，只保留内容
            body_content = ''.join(str(child) for child in body.children)
        else:
            body_content = ""
        
        return styles, body_content
    
    def _scope_css(self, css: str, scope_class: str) -> str:
        """
        为 CSS 添加 scoping，避免样式冲突
        
        例如：
        .title { color: blue; }  =>  .slide-1 .title { color: blue; }
        body { margin: 0; }      =>  .slide-1 { margin: 0; }
        """
        lines = []
        in_media_query = False
        in_keyframes = False
        
        for line in css.split('\n'):
            stripped = line.strip()
            
            # 处理 @media 查询
            if stripped.startswith('@media'):
                in_media_query = True
                lines.append(line)
                continue
            
            # 处理 @keyframes
            if stripped.startswith('@keyframes'):
                in_keyframes = True
                lines.append(line)
                continue
            
            # 检测块结束
            if stripped == '}':
                if in_media_query:
                    in_media_query = False
                elif in_keyframes:
                    in_keyframes = False
                lines.append(line)
                continue
            
            # 如果在 @media 或 @keyframes 中，不添加 scoping
            if in_media_query or in_keyframes:
                lines.append(line)
                continue
            
            # 处理普通样式规则
            if '{' in stripped and not stripped.startswith('@') and not stripped.startswith('/*'):
                selector_part = stripped[:stripped.index('{')].strip()
                rest = stripped[stripped.index('{'):]
                
                # 分割多个选择器（逗号分隔）
                selectors = [s.strip() for s in selector_part.split(',')]
                scoped_selectors = []
                
                for sel in selectors:
                    if not sel:
                        continue
                        
                    # 特殊处理 body 和 html 选择器
                    # 1. 完全匹配 body 或 html -> 替换为 scope_class
                    if sel.lower() in ('body', 'html'):
                        scoped_sel = f".{scope_class}"
                    # 2. body/html 开头的后代选择器 -> 替换开头
                    elif sel.lower().startswith('body ') or sel.lower().startswith('html '):
                        scoped_sel = f".{scope_class} " + sel[5:]
                    # 3. 包含 body/html 的复合选择器 (如 div.body) -> 简单处理，直接前置 scope
                    # 注意：这里简化处理，假设 slide CSS 比较规范
                    else:
                        scoped_sel = f".{scope_class} {sel}"
                        
                    scoped_selectors.append(scoped_sel)
                
                scoped_line = ', '.join(scoped_selectors) + ' ' + rest
                lines.append(' ' * (len(line) - len(line.lstrip())) + scoped_line)
            else:
                lines.append(line)
        
        return '\n'.join(lines)
    
    def _merge_slides_html(self, slides_html: List[str]) -> str:
        """
        合并多个幻灯片的 HTML，直接内联内容（不使用 iframe）

        注意：Playwright 的 page.pdf() 无法正确渲染 iframe 内容，
        因此必须将每个幻灯片的内容直接内联到主文档中。
        """
        all_contents = []
        all_styles = []

        for i, slide_html in enumerate(slides_html):
            # 提取 body 内容
            body_content = self._extract_body_content(slide_html)
            # 提取 style 内容，并添加作用域前缀避免样式冲突
            style_content = self._extract_style_content(slide_html)

            # 为每个幻灯片的样式添加作用域
            scoped_style = self._scope_css(style_content, f"slide-{i+1}")
            all_styles.append(scoped_style)

            # 包装幻灯片内容
            wrapped_content = f'''
            <div class="slide-page slide-{i+1}" id="slide-{i+1}">
                {body_content}
            </div>
            '''
            all_contents.append(wrapped_content)

        # 合并成最终的 HTML
        merged_html = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <link href="https://fonts.googleapis.com/css2?family=Noto+Sans+SC:wght@300;400;500;600;700&display=swap" rel="stylesheet">
    <link href="https://fonts.googleapis.com/icon?family=Material+Icons" rel="stylesheet">
    <style>
        /* 全局样式 */
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        body {{
            background: white;
            font-family: 'Noto Sans SC', 'Microsoft YaHei', sans-serif;
        }}
        .slide-page {{
            width: {self.slide_width}px;
            height: {self.slide_height}px;
            page-break-after: always;
            overflow: hidden;
            position: relative;
        }}
        .slide-page:last-child {{
            page-break-after: auto;
        }}
        /* 幻灯片样式 */
        {chr(10).join(all_styles)}
    </style>
</head>
<body>
    {''.join(all_contents)}
</body>
</html>'''

        return merged_html
        
    async def export_to_pdf(
        self,
        slides_html: List[str],
        title: str = "presentation"
    ) -> Tuple[str, str]:
        """
        导出 PDF - 直接拼接 HTML，不使用 iframe
        
        方案：
        1. 静态化所有幻灯片（Canvas 转图片）
        2. 解析每个幻灯片的 HTML，提取样式和内容
        3. 直接拼接所有幻灯片的内容
        4. 使用 Playwright 渲染并打印为 PDF
        """
        try:
            filename = f"{title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
            filepath = EXPORT_DIR / filename
            
            # 步骤 1：静态化所有幻灯片（将 Canvas 等动态内容转为静态）
            logger.info(f"Staticizing {len(slides_html)} slides for PDF export...")
            slides_html = await batch_staticize_html(slides_html, timeout=30)
            logger.info("All slides staticized")
            
            # 步骤 2：合并所有幻灯片的 HTML
            pdf_html = self._merge_slides_html(slides_html)
            
            async with async_playwright() as p:
                browser = await p.chromium.launch(
                    headless=True,
                    args=['--no-sandbox', '--disable-setuid-sandbox']
                )
                
                page = await browser.new_page()
                
                # 加载 HTML
                await page.set_content(pdf_html, wait_until="domcontentloaded", timeout=120000)
                
                # 等待页面渲染完成（增加等待时间确保 iframe 渲染）
                await page.wait_for_timeout(3000)
                
                # 等待所有 iframe 加载完成
                await page.evaluate("""() => {
                    return Promise.all(
                        Array.from(document.querySelectorAll('iframe')).map(iframe => {
                            return new Promise((resolve) => {
                                if (iframe.contentDocument && iframe.contentDocument.readyState === 'complete') {
                                    resolve();
                                } else {
                                    iframe.addEventListener('load', resolve);
                                    setTimeout(resolve, 2000); // 超时保护
                                }
                            });
                        })
                    );
                }""")
                
                # 再等待一下确保渲染完成
                await page.wait_for_timeout(2000)
                
                # 打印 PDF
                await page.pdf(
                    path=str(filepath),
                    width=f"{self.slide_width}px",
                    height=f"{self.slide_height}px",
                    print_background=True,
                    display_header_footer=False,
                    margin={"top": "0", "right": "0", "bottom": "0", "left": "0"}
                )
                
                await browser.close()
            
            logger.info(f"PDF exported: {filepath}")
            return str(filepath), filename
            
        except Exception as e:
            logger.error(f"Failed to export PDF: {e}")
            raise

    
    def _create_single_slide_html(self, slide_html: str) -> str:
        """为单页渲染创建完整 HTML（包含所有依赖）"""
        body = self._extract_body_content(slide_html)
        style = self._extract_style_content(slide_html)
        
        return f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <script src="https://cdn.tailwindcss.com"></script>
    <script src="https://cdn.jsdelivr.net/npm/echarts@5.4.3/dist/echarts.min.js"></script>
    <link href="https://fonts.googleapis.com/css2?family=Noto+Sans+SC:wght@300;400;500;600;700&display=swap" rel="stylesheet">
    <link href="https://fonts.googleapis.com/icon?family=Material+Icons" rel="stylesheet">
    <style>
        html, body {{
            margin: 0;
            padding: 0;
            width: {self.slide_width}px;
            height: {self.slide_height}px;
            overflow: hidden;
        }}
        * {{ box-sizing: border-box; }}
        {style}
    </style>
</head>
<body>
    {body}
</body>
</html>'''
    
    def _create_screenshot_pdf_html(self, slide_images: List[str]) -> str:
        """创建包含所有截图的 PDF HTML"""
        pages = []
        for i, img_base64 in enumerate(slide_images):
            pages.append(f'''
                <div class="page">
                    <img src="data:image/png;base64,{img_base64}" />
                </div>
            ''')
        
        return f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <style>
        * {{ margin: 0; padding: 0; }}
        .page {{
            width: {self.slide_width}px;
            height: {self.slide_height}px;
            page-break-after: always;
            overflow: hidden;
        }}
        .page:last-child {{
            page-break-after: auto;
        }}
        .page img {{
            width: 100%;
            height: 100%;
            display: block;
        }}
    </style>
</head>
<body>
    {''.join(pages)}
</body>
</html>'''

    async def _prerender_slides(self, browser, slides_html: List[str]) -> List[str]:
        """
        预渲染幻灯片：
        1. 解决 ID 冲突（每页独立渲染）
        2. 将 ECharts 动态图表转换为静态图片/SVG
        3. 确保所有外部资源加载完成
        """
        processed_slides = []
        page = await browser.new_page()
        
        # 预加载 ECharts 脚本，强制使用 SVG 渲染器（如果可能）
        # 注意：这里我们通过篡改 echarts.init 来尝试强制 SVG，如果不生效，后续有 Canvas 转 Image 的兜底
        await page.add_init_script("""
            window.FORCE_SVG = true;
            // 尝试拦截 echarts.init (需要在 echarts 加载后生效，或者由 slide 内部脚本配合)
        """)
        
        for i, slide in enumerate(slides_html):
            try:
                # 构造包含所有依赖的完整 HTML
                temp_html = self._create_temp_slide_html(slide)
                
                # 设置内容并等待加载
                # 使用 domcontentloaded 减少等待时间，因为资源已内联
                # 超时增加到 120s 以应对复杂情况
                await page.set_content(temp_html, wait_until="domcontentloaded", timeout=120000)
                
                # 额外等待动画和图表渲染
                await page.wait_for_timeout(2000)
                
                # 执行客户端脚本：
                # 1. 将 CSSOM 样式固化到 DOM (解决 Tailwind 等动态样式在 DOM 中为空的问题)
                # 2. 将 Canvas 图表转换为图片
                # 3. 移除脚本标签
                # 4. 移除编辑器辅助元素
                await page.evaluate("""() => {
                    // 1. 固化 CSSOM 样式
                    try {
                        for (const sheet of document.styleSheets) {
                            try {
                                // 跳过有 href 的 (link 标签)，只处理内联 style
                                if (!sheet.href && sheet.ownerNode && sheet.ownerNode.tagName === 'STYLE') {
                                    if (sheet.cssRules) {
                                        let css = '';
                                        for (const rule of sheet.cssRules) {
                                            css += rule.cssText + '\\n';
                                        }
                                        if (css) {
                                            sheet.ownerNode.textContent = css;
                                        }
                                    }
                                }
                            } catch (e) {
                                // 忽略跨域样式表访问错误
                            }
                        }
                    } catch (e) {
                        console.error('CSS Materialization failed:', e);
                    }

                    // 2. Canvas -> Image
                    document.querySelectorAll('canvas').forEach(canvas => {
                        try {
                            const img = document.createElement('img');
                            img.src = canvas.toDataURL('image/png');
                            img.style.cssText = canvas.style.cssText;
                            img.className = canvas.className;
                            img.style.width = canvas.width + 'px';
                            img.style.height = canvas.height + 'px';
                            if (canvas.parentNode) {
                                canvas.parentNode.replaceChild(img, canvas);
                            }
                        } catch (e) {
                            console.error('Canvas to Image failed:', e);
                        }
                    });
                    
                    // 3. 移除脚本
                    document.querySelectorAll('script').forEach(el => el.remove());
                    
                    // 4. 移除 contenteditable
                    document.querySelectorAll('[contenteditable]').forEach(el => {
                        el.removeAttribute('contenteditable');
                    });
                }""")
                
                # 获取完整的 HTML 内容 (包含 head 中的 style)
                # 因为只有 body 内容会丢失 <style> 标签定义的样式
                processed_html = await page.content()
                processed_slides.append(processed_html)
                
                logger.info(f"Slide {i+1} pre-rendered successfully")
                
            except Exception as e:
                logger.error(f"Failed to pre-render slide {i+1}: {e}")
                # 如果失败，回退到原始 HTML，避免整个导出失败
                processed_slides.append(slide)
        
        await page.close()
        return processed_slides

    def _create_temp_slide_html(self, slide_content: str) -> str:
        """为单页渲染创建包含完整依赖的 HTML"""
        # 尝试提取 body 内容（如果输入已经是完整 HTML）
        body = self._extract_body_content(slide_content)
        # 提取样式
        style = self._extract_style_content(slide_content)
        
        return f'''
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <script src="https://cdn.tailwindcss.com"></script>
            <script src="https://cdn.jsdelivr.net/npm/echarts@5.4.3/dist/echarts.min.js"></script>
            <link href="https://fonts.googleapis.com/css2?family=Noto+Sans+SC:wght@300;400;500;600;700&display=swap" rel="stylesheet">
            <link href="https://fonts.googleapis.com/icon?family=Material+Icons" rel="stylesheet">
            <style>
                body {{ margin: 0; padding: 0; }}
                * {{ box-sizing: border-box; }}
                {style}
            </style>
        </head>
        <body>
            {body}
        </body>
        </html>
        '''
    
    def _create_pdf_html(self, slides_html: List[str]) -> str:
        """创建用于 PDF 导出的 HTML 文档"""
        slides_content = []
        
        for i, slide_html in enumerate(slides_html):
            # 提取 body 内容
            body_content = self._extract_body_content(slide_html)
            # 提取 style 内容并限定作用域为 .slide-page (避免污染全局和其它页)
            style_content = self._extract_style_content(slide_html, ".slide-page")
            
            slides_content.append(f'''
                <div class="slide-page">
                    <style>{style_content}</style>
                    <div class="slide-inner">
                        {body_content}
                    </div>
                </div>
            ''')
        
        return f'''
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <script src="https://cdn.tailwindcss.com"></script>
            <script src="https://cdn.jsdelivr.net/npm/echarts@5.4.3/dist/echarts.min.js"></script>
            <link href="https://fonts.googleapis.com/css2?family=Noto+Sans+SC:wght@300;400;500;600;700&display=swap" rel="stylesheet">
            <link href="https://fonts.googleapis.com/icon?family=Material+Icons" rel="stylesheet">
            <style>
                body {{
                    margin: 0;
                    padding: 0;
                }}
                .slide-page {{
                    width: 1280px;
                    height: 720px;
                    page-break-after: always;
                    overflow: hidden;
                    position: relative;
                }}
                .slide-inner {{
                    width: 100%;
                    height: 100%;
                    position: relative;
                }}
                * {{
                    box-sizing: border-box;
                    font-family: 'Noto Sans SC', 'Microsoft YaHei', sans-serif;
                }}
            </style>
        </head>
        <body>
            {''.join(slides_content)}
        </body>
        </html>
        '''
    
    def _extract_body_content(self, html: str) -> str:
        """从 HTML 中提取 body 内容"""
        import re
        body_match = re.search(r'<body[^>]*>(.*?)</body>', html, re.DOTALL | re.IGNORECASE)
        if body_match:
            return body_match.group(1)
        return html
    
    def _extract_style_content(self, html: str, scope_selector: str = None) -> str:
        """
        从 HTML 中提取 style 内容
        
        Args:
            html: HTML 内容
            scope_selector: CSS 作用域选择器 (例如 .slide-page)，如果提供，将把 body/html 选择器替换为此选择器
        """
        soup = BeautifulSoup(html, 'html.parser')
        styles = []
        
        for style_tag in soup.find_all('style'):
            if not style_tag.string:
                continue
                
            css_text = style_tag.string
            
            # 过滤掉编辑器特定的重置样式
            if '[contenteditable="true"]' in css_text:
                continue
            if 'width: 1280px !important' in css_text.replace(' ', ''): # 简单去空匹配
                continue
            
            if scope_selector:
                # 替换 html, body 为 scope_selector
                # 1. 处理 "html, body" 组合
                css_text = re.sub(r'html\s*,\s*body', scope_selector, css_text, flags=re.IGNORECASE)
                # 2. 处理单独的 body
                css_text = re.sub(r'(?<![\w-])body(?![\w-])', scope_selector, css_text, flags=re.IGNORECASE)
                # 3. 处理单独的 html (通常移除或替换)
                css_text = re.sub(r'(?<![\w-])html(?![\w-])', scope_selector, css_text, flags=re.IGNORECASE)
                
            styles.append(css_text)
            
        style_content = '\n'.join(styles)
        
        # 替换字体为 PDF 导出友好的字体
        style_content = self._replace_fonts_for_pdf(style_content)
        
        return style_content
    
    def _replace_fonts_for_pdf(self, css: str) -> str:
        """替换字体为 PDF 导出友好的字体"""
        for web_font, pdf_font in self.FONT_REPLACEMENTS.items():
            css = css.replace(web_font, pdf_font)
        
        # 移除 @import 语句（WeasyPrint 可能无法加载）
        import re
        css = re.sub(r'@import\s+url\([^)]+\);?', '', css)
        
        return css
    
    async def export_to_html(
        self, 
        slides_html: List[str], 
        title: str = "presentation"
    ) -> Tuple[str, str]:
        """
        将 HTML 幻灯片导出为单一 HTML 文件
        
        方案：使用 iframe 隔离每个幻灯片的样式，避免 CSS 冲突
        """
        import html as html_module
        
        try:
            filename = f"{title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
            filepath = EXPORT_DIR / filename
            
            # 构建 iframe 版本的 HTML
            slides_content = []
            for i, slide_html in enumerate(slides_html):
                # 确保幻灯片 HTML 包含必要的依赖
                full_slide_html = self._ensure_slide_dependencies(slide_html)
                # 转义 HTML 用于 srcdoc 属性
                escaped_html = html_module.escape(full_slide_html)
                
                slides_content.append(f'''
                    <div class="slide-wrapper" id="slide-{i+1}">
                        <iframe 
                            srcdoc="{escaped_html}"
                            class="slide-frame"
                            frameborder="0"
                            scrolling="no"
                            sandbox="allow-scripts allow-same-origin"
                        ></iframe>
                        <div class="slide-number">{i+1} / {len(slides_html)}</div>
                    </div>
                ''')
            
            full_html = f'''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title} - 演示文稿</title>
    <style>
        * {{
            box-sizing: border-box;
            margin: 0;
            padding: 0;
        }}
        body {{
            font-family: 'Microsoft YaHei', sans-serif;
            background-color: #1a1a2e;
            padding: 40px 20px;
            display: flex;
            flex-direction: column;
            align-items: center;
            gap: 30px;
        }}
        .slide-wrapper {{
            position: relative;
            width: {self.slide_width}px;
            height: {self.slide_height}px;
            background: white;
            box-shadow: 0 8px 32px rgba(0, 0, 0, 0.3);
            border-radius: 8px;
            overflow: hidden;
            flex-shrink: 0;
        }}
        .slide-frame {{
            width: 100%;
            height: 100%;
            border: none;
            display: block;
        }}
        .slide-number {{
            position: absolute;
            bottom: 10px;
            right: 20px;
            background: rgba(0, 0, 0, 0.6);
            color: white;
            padding: 4px 12px;
            border-radius: 4px;
            font-size: 12px;
            pointer-events: none;
            opacity: 0;
            transition: opacity 0.3s;
        }}
        .slide-wrapper:hover .slide-number {{
            opacity: 1;
        }}
        @media print {{
            body {{
                background: none;
                padding: 0;
                gap: 0;
            }}
            .slide-wrapper {{
                box-shadow: none;
                border-radius: 0;
                page-break-after: always;
            }}
            .slide-number {{
                display: none;
            }}
        }}
    </style>
</head>
<body>
    {''.join(slides_content)}
</body>
</html>'''
            
            with open(filepath, "w", encoding="utf-8") as f:
                f.write(full_html)
            
            logger.info(f"HTML exported: {filepath}")
            return str(filepath), filename
            
        except Exception as e:
            logger.error(f"Failed to export HTML: {e}")
            raise
    
    def _ensure_slide_dependencies(self, slide_html: str) -> str:
        """确保幻灯片 HTML 包含必要的依赖（Tailwind, ECharts, 字体等）"""
        # 检查是否已经有 head 标签
        if '<head>' in slide_html.lower():
            # 在 head 中注入依赖
            deps = '''
    <script src="https://cdn.tailwindcss.com"></script>
    <script src="https://cdn.jsdelivr.net/npm/echarts@5.4.3/dist/echarts.min.js"></script>
    <link href="https://fonts.googleapis.com/css2?family=Noto+Sans+SC:wght@300;400;500;600;700&display=swap" rel="stylesheet">
    <link href="https://fonts.googleapis.com/icon?family=Material+Icons" rel="stylesheet">
'''
            return slide_html.replace('</head>', deps + '</head>')
        else:
            # 包装成完整的 HTML
            body = self._extract_body_content(slide_html)
            style = self._extract_style_content(slide_html)
            
            return f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <script src="https://cdn.tailwindcss.com"></script>
    <script src="https://cdn.jsdelivr.net/npm/echarts@5.4.3/dist/echarts.min.js"></script>
    <link href="https://fonts.googleapis.com/css2?family=Noto+Sans+SC:wght@300;400;500;600;700&display=swap" rel="stylesheet">
    <link href="https://fonts.googleapis.com/icon?family=Material+Icons" rel="stylesheet">
    <style>
        html, body {{
            margin: 0;
            padding: 0;
            width: 100%;
            height: 100%;
            overflow: hidden;
        }}
        {style}
    </style>
</head>
<body>
    {body}
</body>
</html>'''

    def _create_standalone_html(self, slides_html: List[str], title: str) -> str:
        """创建独立的 HTML 浏览文件"""
        slides_content = []
        
        for i, slide_html in enumerate(slides_html):
            # 提取内容
            body_content = self._extract_body_content(slide_html)
            # 提取内容
            body_content = self._extract_body_content(slide_html)
            # 提取 style 内容并限定作用域为 .slide-content
            style_content = self._extract_style_content(slide_html, ".slide-content")
            
            slides_content.append(f'''
                <!-- Slide {i+1} -->
                <div class="slide-wrapper" id="slide-{i+1}">
                    <div class="slide-content">
                        <style>{style_content}</style>
                        {body_content}
                    </div>
                    <div class="slide-number">{i+1} / {len(slides_html)}</div>
                </div>
            ''')
        
        return f'''
        <!DOCTYPE html>
        <html lang="zh-CN">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>{title} - 演示文稿</title>
            <script src="https://cdn.tailwindcss.com"></script>
            <script src="https://cdn.jsdelivr.net/npm/echarts@5.4.3/dist/echarts.min.js"></script>
            <link href="https://fonts.googleapis.com/css2?family=Noto+Sans+SC:wght@300;400;500;600;700&display=swap" rel="stylesheet">
            <link href="https://fonts.googleapis.com/icon?family=Material+Icons" rel="stylesheet">
            <style>
                * {{
                    box-sizing: border-box;
                    margin: 0;
                    padding: 0;
                }}
                body {{
                    font-family: 'Noto Sans SC', sans-serif;
                    background-color: #f0f2f5;
                    padding: 40px 20px;
                    display: flex;
                    flex-direction: column;
                    align-items: center;
                    gap: 30px;
                }}
                .slide-wrapper {{
                    position: relative;
                    width: 1280px;
                    height: 720px;
                    background: white;
                    box-shadow: 0 4px 12px rgba(0, 0, 0, 0.1);
                    border-radius: 8px;
                    overflow: hidden;
                    flex-shrink: 0;
                }}
                .slide-content {{
                    width: 100%;
                    height: 100%;
                    position: relative;
                    overflow: hidden;
                }}
                /* 确保幻灯片内容占满容器 */
                .slide-content > div {{
                    width: 100%;
                    height: 100%;
                }}
                .slide-number {{
                    position: absolute;
                    bottom: 10px;
                    right: 20px;
                    background: rgba(0, 0, 0, 0.5);
                    color: white;
                    padding: 4px 8px;
                    border-radius: 4px;
                    font-size: 12px;
                    pointer-events: none;
                    opacity: 0;
                    transition: opacity 0.3s;
                }}
                .slide-wrapper:hover .slide-number {{
                    opacity: 1;
                }}
                /* 滚动条美化 */
                ::-webkit-scrollbar {{
                    width: 10px;
                    background: transparent;
                }}
                ::-webkit-scrollbar-thumb {{
                    background: #ccc;
                    border-radius: 5px;
                }}
                
                @media print {{
                    body {{
                        background: none;
                        display: block;
                        padding: 0;
                    }}
                    .slide-wrapper {{
                        box-shadow: none;
                        margin: 0;
                        page-break-after: always;
                        border-radius: 0;
                    }}
                    .slide-number {{
                        display: none;
                    }}
                }}
            </style>
        </head>
        <body>
            {''.join(slides_content)}
        </body>
        </html>
        '''
    
    async def export_to_images(
        self, 
        slides_html: List[str], 
        title: str = "presentation",
        format: str = "png"
    ) -> Tuple[str, str]:
        """
        将 HTML 幻灯片导出为图片（打包为 ZIP）
        
        Args:
            slides_html: HTML 幻灯片列表
            title: 文件标题
            format: 图片格式 (png/jpg)
            
        Returns:
            (文件路径, 文件名)
        """
        try:
            # 静态化所有幻灯片
            logger.info(f"Staticizing {len(slides_html)} slides for image export...")
            slides_html = await batch_staticize_html(slides_html, timeout=30)
            logger.info("All slides staticized")
            
            # 生成唯一文件名
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            zip_filename = f"{title}_{timestamp}_images.zip"
            zip_filepath = EXPORT_DIR / zip_filename
            
            # 使用 Playwright 截图
            async with async_playwright() as p:
                browser = await p.chromium.launch(
                    headless=True,
                    args=['--no-sandbox', '--disable-setuid-sandbox']
                )
                page = await browser.new_page(viewport={'width': self.slide_width, 'height': self.slide_height})
                
                # 创建 ZIP 文件
                with zipfile.ZipFile(zip_filepath, 'w', zipfile.ZIP_DEFLATED) as zipf:
                    for i, slide_html in enumerate(slides_html):
                        # 设置页面内容
                        await page.set_content(slide_html)
                        await page.wait_for_load_state('networkidle')
                        
                        # 额外等待字体和图片加载
                        await asyncio.sleep(0.5)
                        
                        # 等待所有图片加载完成
                        await page.evaluate('''
                            () => {
                                return Promise.all(
                                    Array.from(document.images)
                                        .filter(img => !img.complete)
                                        .map(img => new Promise(resolve => {
                                            img.onload = img.onerror = resolve;
                                        }))
                                );
                            }
                        ''')
                        
                        # 截图
                        screenshot = await page.screenshot(type=format, full_page=False)
                        
                        # 添加到 ZIP
                        image_filename = f"slide_{i+1:02d}.{format}"
                        zipf.writestr(image_filename, screenshot)
                        
                        logger.info(f"Screenshot captured: {image_filename}")
                
                await browser.close()
            
            logger.info(f"Images exported: {zip_filepath}")
            return str(zip_filepath), zip_filename
            
        except Exception as e:
            logger.error(f"Failed to export images: {e}")
            raise
    
    async def export_to_pptx(
        self, 
        slides_html: List[str], 
        title: str = "presentation"
    ) -> Tuple[str, str]:
        """
        将 HTML 幻灯片导出为可编辑的 PPTX 文件
        
        技术方案（参考智谱清言）：
        1. 使用 Playwright 渲染 HTML，提取元素位置和样式
        2. 使用 python-pptx 创建 PPTX 文件
        3. 将每个元素转换为对应的 PPTX 元素（文本框、形状、图片）
        
        优点：
        - 文本可编辑
        - 形状可编辑
        - 保持原始布局
        
        Args:
            slides_html: HTML 幻灯片列表
            title: 文件标题
            
        Returns:
            (文件路径, 文件名)
        """
        try:
            # 生成唯一文件名
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            pptx_filename = f"{title}_{timestamp}.pptx"
            pptx_filepath = EXPORT_DIR / pptx_filename
            
            # 使用新的 PPTX 生成器（HTML 渲染 + 元素截图 + python-pptx）
            await convert_html_to_pptx(
                slides_html=slides_html,
                output_path=str(pptx_filepath),
                title=title
            )
            
            # 验证文件是否生成成功
            if not pptx_filepath.exists():
                raise Exception(f"PPTX 文件未生成: {pptx_filepath}")
            
            # 验证文件大小（确保不是空文件）
            file_size = pptx_filepath.stat().st_size
            if file_size < 1000:  # 小于 1KB 可能是损坏的文件
                raise Exception(f"PPTX 文件可能损坏，大小只有 {file_size} 字节")
            
            logger.info(f"PPTX exported successfully: {pptx_filepath} ({file_size} bytes)")
            return str(pptx_filepath), pptx_filename
            
        except Exception as e:
            logger.error(f"Failed to export PPTX: {e}")
            raise


# 单例实例
exporter = PPTExporter()


async def pre_export_check(slides_html: List[str]) -> Dict[str, Any]:
    """
    导出前预检
    
    Args:
        slides_html: HTML 幻灯片列表
        
    Returns:
        预检结果 {
            "passed": bool,
            "slides_count": int,
            "issues": List[str],
            "warnings": List[str],
            "external_resources_count": int,
            "external_resources": List[str]
        }
    """
    all_issues = []
    all_warnings = []
    all_resources = []
    
    for i, html in enumerate(slides_html):
        result = validate_html(html)
        
        if result["issues"]:
            all_issues.extend([f"第{i+1}页: {issue}" for issue in result["issues"]])
        
        if result["warnings"]:
            all_warnings.extend([f"第{i+1}页: {warning}" for warning in result["warnings"]])
        
        all_resources.extend(result["external_resources"])
    
    return {
        "passed": len(all_issues) == 0,
        "slides_count": len(slides_html),
        "issues": all_issues,
        "warnings": all_warnings,
        "external_resources_count": len(set(all_resources)),
        "external_resources": list(set(all_resources))
    }


async def generate_export_report(
    slides_html: List[str],
    export_format: str,
    output_path: str,
    pre_check_result: Optional[Dict[str, Any]] = None
) -> Dict[str, Any]:
    """
    生成导出质量报告
    
    Args:
        slides_html: HTML 幻灯片列表
        export_format: 导出格式
        output_path: 输出文件路径
        pre_check_result: 预检结果
        
    Returns:
        质量报告
    """
    report = {
        "format": export_format,
        "output_path": output_path,
        "slides_count": len(slides_html),
        "export_time": datetime.now().isoformat(),
        "file_exists": os.path.exists(output_path),
        "file_size": 0,
        "issues": [],
        "warnings": [],
        "quality_score": 1.0
    }
    
    # 检查文件是否存在
    if not report["file_exists"]:
        report["issues"].append("导出文件不存在")
        report["quality_score"] = 0
        return report
    
    # 获取文件大小
    report["file_size"] = os.path.getsize(output_path)
    
    # 检查文件大小
    if report["file_size"] < 1000:
        report["issues"].append(f"文件大小异常（{report['file_size']} 字节）")
        report["quality_score"] -= 0.5
    
    # 检查 PPTX 文件完整性
    if export_format == "pptx":
        try:
            from pptx import Presentation
            prs = Presentation(output_path)
            report["slides_in_pptx"] = len(prs.slides)
            
            if report["slides_in_pptx"] != report["slides_count"]:
                report["warnings"].append(
                    f"PPTX 幻灯片数量不匹配：期望 {report['slides_count']}，实际 {report['slides_in_pptx']}"
                )
                report["quality_score"] -= 0.2
            
            # 检查每页是否有内容
            empty_slides = []
            for i, slide in enumerate(prs.slides):
                if len(slide.shapes) == 0:
                    empty_slides.append(i + 1)
            
            if empty_slides:
                report["warnings"].append(f"第 {', '.join(map(str, empty_slides))} 页没有内容")
                report["quality_score"] -= 0.1 * len(empty_slides)
        
        except Exception as e:
            report["issues"].append(f"PPTX 文件可能损坏: {e}")
            report["quality_score"] -= 0.5
    
    # 添加预检结果
    if pre_check_result:
        if pre_check_result.get("issues"):
            report["warnings"].extend(pre_check_result["issues"])
        if pre_check_result.get("warnings"):
            report["warnings"].extend(pre_check_result["warnings"])
    
    # 确保质量分数在 0-1 之间
    report["quality_score"] = max(0, min(1, report["quality_score"]))
    
    return report


async def export_ppt(
    slides_html: List[str],
    format: str,
    title: str = "presentation",
    skip_validation: bool = False
) -> Tuple[str, str]:
    """
    导出 PPT（增加预检功能）
    
    Args:
        slides_html: HTML 幻灯片列表
        format: 导出格式 (pdf/png/pptx)
        title: 文件标题
        skip_validation: 是否跳过验证
        
    Returns:
        (文件路径, 文件名)
    """
    # 导出前预检
    if not skip_validation:
        check_result = await pre_export_check(slides_html)
        
        if not check_result["passed"]:
            logger.warning(f"导出预检发现问题: {check_result['issues']}")
            # 可以选择抛出异常或继续导出
            # raise ValueError(f"导出预检失败: {check_result['issues']}")
        
        if check_result["warnings"]:
            logger.warning(f"导出预检警告: {check_result['warnings']}")
    
    # 原有的导出逻辑
    if format == "pdf":
        return await exporter.export_to_pdf(slides_html, title)
    elif format == "html":
        return await exporter.export_to_html(slides_html, title)
    elif format in ["png", "jpg", "images"]:
        return await exporter.export_to_images(slides_html, title, "png")
    elif format == "pptx":
        return await exporter.export_to_pptx(slides_html, title)
    else:
        raise ValueError(f"Unsupported format: {format}")
