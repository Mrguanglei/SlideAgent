
import os
import json
import logging
import asyncio
import base64
import time
import io
from pathlib import Path
from typing import List, Dict, Optional, Tuple, Any

import google.generativeai as genai
from pdf2image import convert_from_path
from PIL import Image, ImageFile
from pptx import Presentation

# Fix for "image file is truncated" errors
ImageFile.LOAD_TRUNCATED_IMAGES = True
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

class ExportToolService:
    def __init__(self):
        self.google_api_base = os.getenv("GOOGLE_API_BASE")
        
        if not self.api_key:
            logger.warning("GOOGLE_API_KEY not found. ExportToolService will fail if used.")
        else:
            # Configure Gemini
            # If using a custom base URL (proxy), we often need REST transport
            if self.google_api_base:
                from google.api_core.client_options import ClientOptions
                logger.info(f"Using Custom Google API Base: {self.google_api_base}")
                genai.configure(
                    api_key=self.api_key,
                    transport="rest",
                    client_options=ClientOptions(api_endpoint=self.google_api_base)
                )
            else:
                genai.configure(api_key=self.api_key)
        # Models configurations
        # Note: image model availability varies. 
        # Ideally: gemini-2.5-flash-image (if available) or gemini-2.5-flash
        self.model_text_removal = "gemini-2.5-flash" 
        self.model_ocr = "gemini-2.5-flash-lite" # Cost effective

    async def process_pdf(self, pdf_path: str, output_path: str) -> str:
        """
        Convert a PDF file to an editable PPTX file.
        Returns the path to the generated PPTX.
        """
        if not self.api_key:
            raise ValueError("GOOGLE_API_KEY is required forExportToolService")
        
        logger.info(f"Processing PDF: {pdf_path}")
        
        # 1. Convert PDF to images
        images = await self._pdf_to_images(pdf_path)
        logger.info(f"Converted PDF to {len(images)} images")
        
        slides_data = []
        
        # 2. Process each page (Parallelize could be better, but sequential for stability first)
        for i, img in enumerate(images):
            try:
                logger.info(f"Processing page {i+1}/{len(images)}")
                
                # Run Text Removal and OCR in parallel
                cleaned_img_task = self._remove_text(img)
                ocr_data_task = self._extract_text(img)
                
                cleaned_img, ocr_data = await asyncio.gather(cleaned_img_task, ocr_data_task)
                
                slides_data.append({
                    "original_size": img.size,
                    "cleaned_image": cleaned_img,
                    "text_blocks": ocr_data
                })
                
            except Exception as e:
                logger.error(f"Error processing page {i+1}: {e}")
                # Fallback: use original image and empty text
                slides_data.append({
                    "original_size": img.size,
                    "cleaned_image": img, # Use original if cleaning fails
                    "text_blocks": []
                })

        # 3. Create PPTX
        self._create_pptx(slides_data, output_path)
        logger.info(f"PPTX saved to {output_path}")
        
        return output_path

    async def _pdf_to_images(self, pdf_path: str) -> List[Image.Image]:
        """Convert PDF to a list of PIL Images using pdf2image"""
        # Run in executor to avoid blocking event loop
        loop = asyncio.get_event_loop()
        images = await loop.run_in_executor(None, convert_from_path, pdf_path, 200) # 200 DPI
        # Force load images to prevent "truncated" errors with Gemini SDK
        for img in images:
            img.load()
        return images

    async def _remove_text(self, image: Image.Image) -> Image.Image:
        """Call Gemini to remove text from image"""
        try:
            model = genai.GenerativeModel(self.model_text_removal)
            
            # More explicit prompt for standard models if they support it
            prompt = "Return a copy of this image with ALL text removed. Preserve the background exactly."
            
            # Run in executor
            loop = asyncio.get_event_loop()
            response = await loop.run_in_executor(None, model.generate_content, [prompt, image])
            
            # Try to get the image from response
            # Standard SDK way to check for image parts
            if hasattr(response, 'parts'):
                for part in response.parts:
                    # Check if part is an image (SDK dependent, but try standard attribute access)
                    # For some versions, it might be inline_data
                    if hasattr(part, 'inline_data') and part.inline_data:
                        img_data = base64.b64decode(part.inline_data.data)
                        return Image.open(io.BytesIO(img_data))
                    # Or if it's a specific 'image' type part
                    if hasattr(part, 'image'):
                         return part.image
            
            # If we get here, SDK didn't give us an image easily.
            # Fallback for now to avoid breaking the flow.
            # Ideally we would use the REST API directly for consistent behavior across versions.
            logger.warning("Gemini did not return an image part. Returning original.")
            return image

        except Exception as e:
            logger.warning(f"Text removal failed: {e}")
            return image

    async def _extract_text(self, image: Image.Image) -> List[Dict]:
        """OCR to extract text blocks with styling"""
        model = genai.GenerativeModel(self.model_ocr)
        
        prompt = """
        Analyze this image and extract all text blocks with precise positioning and styling.
        For each text block, return a JSON object with:
        - text: the exact text content
        - box_2d: bounding box as [ymin, xmin, ymax, xmax] in 0-1000 coordinate system
        - font_size_pt: estimated font size in points (typical range: 8-72)
        - font_weight: "normal" or "bold"
        - font_style: "normal" or "italic"
        - text_align: "left", "center", or "right"
        - color: hex color code like "000000" or "FFFFFF"
        
        Return a JSON ARRAY.
        """
        
        try:
            loop = asyncio.get_event_loop()
            response = await loop.run_in_executor(
                None, 
                lambda: model.generate_content(
                    [prompt, image], 
                    generation_config={"response_mime_type": "application/json"}
                )
            )
            
            text_result = response.text
            if not text_result:
                return []
                
            return json.loads(text_result)
            
        except Exception as e:
            logger.warning(f"OCR failed: {e}")
            return []

    def _create_pptx(self, slides_data: List[Dict], output_path: str):
        """Generate PPTX from processed data"""
        prs = Presentation()
        
        # Set slide size to 16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        blank_layout = prs.slide_layouts[6] # 6 = Blank
        
        for slide_data in slides_data:
            slide = prs.slides.add_slide(blank_layout)
            
            # 1. Add background image
            # We need to save PIL image to bytes/tempfile to add to PPTX
            img_stream = split_image_to_stream(slide_data["cleaned_image"])
            slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            # 2. Add text blocks
            for block in slide_data["text_blocks"]:
                self._add_text_block(slide, block, prs.slide_width, prs.slide_height)
                
        prs.save(output_path)

    def _add_text_block(self, slide, block, slide_w, slide_h):
        try:
            box = block.get("box_2d")
            if not box or len(box) != 4:
                return

            ymin, xmin, ymax, xmax = box
            
            # Convert 0-1000 coords to EMU
            left = int((xmin / 1000) * slide_w)
            top = int((ymin / 1000) * slide_h)
            width = int(((xmax - xmin) / 1000) * slide_w)
            height = int(((ymax - ymin) / 1000) * slide_h)
            
            shape = slide.shapes.add_textbox(left, top, width, height)
            tf = shape.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            text_content = block.get("text", "")
            p.text = text_content
            
            # Align
            align_map = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT
            }
            p.alignment = align_map.get(block.get("text_align", "left"), PP_ALIGN.LEFT)
            
            # Font style
            run = p.runs[0]
            font = run.font
            
            # Intelligent font selection
            # If text contains non-ASCII (likely Chinese), use generic CJK font
            if any(ord(c) > 127 for c in text_content):
                font.name = 'Microsoft YaHei' 
            else:
                font.name = 'Arial'
            
            sz = block.get("font_size_pt")
            if sz:
                font.size = Pt(float(sz))
            
            if block.get("font_weight") == "bold":
                font.bold = True
            
            if block.get("font_style") == "italic":
                font.italic = True
                
            color_hex = block.get("color", "000000").replace("#", "")
            if len(color_hex) == 6:
                try:
                    font.color.rgb = RGBColor.from_string(color_hex)
                except:
                    pass
                    
        except Exception as e:
            logger.warning(f"Error adding text block: {e}")

def split_image_to_stream(image: Image.Image):
    import io
    stream = io.BytesIO()
    image.save(stream, format="PNG")
    stream.seek(0)
    return stream

export_tool_service=ExportToolService()